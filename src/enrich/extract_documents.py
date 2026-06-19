"""Extract text from PDF, Excel, Word, PowerPoint, HTML, plain text, iCal, vCard, Numbers, RTF, and AI attachments."""

import json
import os
import subprocess
from pathlib import Path

DATA_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "data")
INPUT_JSONL = os.path.join(DATA_DIR, "2_enrich", "b_pictures.jsonl")
ATTACHMENTS_DIR = os.path.join(DATA_DIR, "1_parser", "attachments")
OUTPUT_JSONL = os.path.join(DATA_DIR, "2_enrich", "c_documents.jsonl")

PDF_EXTS = {".pdf"}
EXCEL_EXTS = {".xlsx", ".xls"}
WORD_EXTS = {".docx"}
PPTX_EXTS = {".pptx"}
HTML_EXTS = {".htm", ".html"}
TEXT_EXTS = {".txt"}
LEGACY_WORD_EXTS = {".doc"}
CAL_EXTS = {".ics"}
CONTACT_EXTS = {".vcf"}
NUMBERS_EXTS = {".numbers"}
RTF_EXTS = {".rtf"}
AI_EXTS = {".ai"}
ALL_SUPPORTED = (
    PDF_EXTS | EXCEL_EXTS | WORD_EXTS | PPTX_EXTS
    | HTML_EXTS | TEXT_EXTS | LEGACY_WORD_EXTS | CAL_EXTS | CONTACT_EXTS
    | NUMBERS_EXTS | RTF_EXTS | AI_EXTS
)


def extract_pdf(path: Path) -> str:
    import pymupdf

    text_parts = []
    with pymupdf.open(str(path)) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts).strip()


def extract_excel(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts).strip()


def extract_word(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs).strip()


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text)
    return "\n".join(parts).strip()


def extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    raw = path.read_bytes()
    # Try utf-8 first, fall back to latin-1
    try:
        html = raw.decode("utf-8")
    except UnicodeDecodeError:
        html = raw.decode("latin-1", errors="ignore")

    soup = BeautifulSoup(html, "html.parser")
    # Remove script and style elements
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def extract_plain_text(path: Path) -> str:
    raw = path.read_bytes()
    try:
        return raw.decode("utf-8").strip()
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="ignore").strip()


def extract_legacy_word(path: Path) -> str:
    """Extract text from legacy .doc files using antiword if available, else raw extraction."""
    try:
        result = subprocess.run(
            ["antiword", str(path)],
            capture_output=True, text=True, check=True, timeout=30,
        )
        return result.stdout.strip()
    except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
        pass

    # Fallback: extract printable text from binary
    raw = path.read_bytes()
    # Extract runs of printable ASCII (at least 4 chars long)
    import re
    text_runs = re.findall(rb"[\x20-\x7e]{4,}", raw)
    return "\n".join(r.decode("ascii") for r in text_runs).strip()


def extract_ical(path: Path) -> str:
    from icalendar import Calendar

    raw = path.read_bytes()
    cal = Calendar.from_ical(raw)
    parts = []
    for component in cal.walk():
        if component.name == "VEVENT":
            summary = str(component.get("SUMMARY", ""))
            dtstart = component.get("DTSTART")
            dtend = component.get("DTEND")
            location = str(component.get("LOCATION", ""))
            description = str(component.get("DESCRIPTION", ""))

            event_parts = [f"Event: {summary}"]
            if dtstart:
                event_parts.append(f"Start: {dtstart.dt}")
            if dtend:
                event_parts.append(f"End: {dtend.dt}")
            if location:
                event_parts.append(f"Location: {location}")
            if description:
                event_parts.append(f"Description: {description}")
            parts.append("\n".join(event_parts))
    return "\n\n".join(parts).strip()


def extract_vcard(path: Path) -> str:
    import vobject

    raw = path.read_bytes()
    try:
        text = raw.decode("utf-8")
    except UnicodeDecodeError:
        text = raw.decode("latin-1", errors="ignore")

    parts = []
    for vcard in vobject.readComponents(text):
        contact_parts = []
        if hasattr(vcard, "fn"):
            contact_parts.append(f"Name: {vcard.fn.value}")
        if hasattr(vcard, "org"):
            org_val = vcard.org.value
            if isinstance(org_val, list):
                org_val = ", ".join(str(o) for o in org_val if o)
            contact_parts.append(f"Organization: {org_val}")
        if hasattr(vcard, "tel"):
            contact_parts.append(f"Phone: {vcard.tel.value}")
        if hasattr(vcard, "email"):
            contact_parts.append(f"Email: {vcard.email.value}")
        if contact_parts:
            parts.append("\n".join(contact_parts))
    return "\n\n".join(parts).strip()


def extract_numbers(path: Path) -> str:
    """Extract text from Apple Numbers spreadsheets."""
    from numbers_parser import Document

    doc = Document(str(path))
    parts = []
    for sheet in doc.sheets:
        parts.append(f"[Sheet: {sheet.name}]")
        for table in sheet.tables:
            parts.append(f"  [Table: {table.name}]")
            for row in range(table.num_rows):
                cells = []
                for col in range(table.num_cols):
                    val = table.cell(row, col).value
                    cells.append(str(val) if val is not None else "")
                parts.append("  " + " | ".join(cells))
    return "\n".join(parts).strip()


def extract_rtf(path: Path) -> str:
    """Extract plain text from RTF files."""
    from striprtf.striprtf import rtf_to_text

    raw = path.read_bytes()
    try:
        text = raw.decode("utf-8")
    except UnicodeDecodeError:
        text = raw.decode("latin-1", errors="ignore")
    return rtf_to_text(text).strip()


def extract_text_content(path: Path) -> str | None:
    ext = path.suffix.lower()
    try:
        if ext in PDF_EXTS or ext in AI_EXTS:
            return extract_pdf(path)
        elif ext in EXCEL_EXTS:
            return extract_excel(path)
        elif ext in WORD_EXTS:
            return extract_word(path)
        elif ext in PPTX_EXTS:
            return extract_pptx(path)
        elif ext in HTML_EXTS:
            return extract_html(path)
        elif ext in TEXT_EXTS:
            return extract_plain_text(path)
        elif ext in LEGACY_WORD_EXTS:
            return extract_legacy_word(path)
        elif ext in CAL_EXTS:
            return extract_ical(path)
        elif ext in CONTACT_EXTS:
            return extract_vcard(path)
        elif ext in NUMBERS_EXTS:
            return extract_numbers(path)
        elif ext in RTF_EXTS:
            return extract_rtf(path)
    except Exception as e:
        print(f"  Warning: Failed to extract {path.name}: {e}")
    return None


def main():
    os.makedirs(os.path.dirname(OUTPUT_JSONL), exist_ok=True)

    with open(INPUT_JSONL, "r", encoding="utf-8") as inp, \
         open(OUTPUT_JSONL, "w", encoding="utf-8") as out:

        lines = inp.readlines()
        total = len(lines)

        for idx, line in enumerate(lines, 1):
            data = json.loads(line)

            for attach in data.get("attachments", []):
                att_path = Path(ATTACHMENTS_DIR) / attach.get("path", "")
                if att_path.suffix.lower() not in ALL_SUPPORTED:
                    continue
                text = extract_text_content(att_path)
                if text:
                    attach["extracted_text"] = text
                    print(f"  Extracted {att_path.name} ({len(text)} chars)")

            out.write(json.dumps(data, ensure_ascii=False) + "\n")
            print(f"[{idx}/{total}] {data['id']}")


if __name__ == "__main__":
    main()
